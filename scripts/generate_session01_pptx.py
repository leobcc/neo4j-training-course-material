"""Generate Session 1 PPTX using Presentation_template.pptx as base.
Italian, Neo4j-centric. Logo on content slides (right). No project cards.
Swirl on title slide only."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree
from pptx.oxml.ns import qn
import os

# ── Paths ──────────────────────────────────────────────────────
TEMPLATE  = r'C:\Users\leob3\OneDrive\Desktop\code\neo4j\Presentation_template.pptx'
ASSETS    = r'C:\Users\leob3\OneDrive\Desktop\code\neo4j\slides\session-01'
SWIRL_PATH = os.path.join(ASSETS, 'swirl.png')

# ── Colours ────────────────────────────────────────────────────
DARK_BG    = RGBColor(0x27, 0x29, 0x36)  # all slides (same bg)
MED_BLUE   = RGBColor(0x1B, 0x2A, 0x4A)
LIGHT_BLUE = RGBColor(0x2D, 0x4F, 0x7A)
GOLD       = RGBColor(0xF0, 0xC0, 0x40)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
CODE_GREEN = RGBColor(0x7E, 0xC6, 0x8E)
ORANGE     = RGBColor(0xF0, 0x80, 0x40)

# ── Layout constants ───────────────────────────────────────────
MARGIN     = Inches(0.8)
CONTENT_W  = Inches(9.4)
CODE_W     = Inches(9.4)
LINE_W     = Inches(11.7)

# ── Helpers ────────────────────────────────────────────────────
def _set_bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color

def _shape(slide, l, t, w, h, fill=None, line=None, line_w=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.line.fill.background()
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line:
        s.line.color.rgb = line; s.line.width = line_w or Pt(1)
    return s

def _tb(slide, l, t, w, h, text, sz=18, color=WHITE, bold=False,
        align=PP_ALIGN.LEFT, font="Calibri"):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text; p.font.size = Pt(sz); p.font.color.rgb = color
    p.font.bold = bold; p.font.name = font; p.alignment = align
    return tb

def _title_heading(slide, text):
    _tb(slide, MARGIN, Inches(0.35), CONTENT_W, Inches(0.65),
        text, sz=33, color=GOLD, bold=True)

def _bullets(slide, l, t, w, h, items, sz=15, color=WHITE, spacing=1.3):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"\u2022  {item}"
        p.font.size = Pt(sz); p.font.color.rgb = color
        p.font.name = "Calibri"; p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(int(sz * (spacing - 1) * 3))

def _code(slide, l, t, w, h, lines, sz=11):
    _shape(slide, l, t, w, h, fill=MED_BLUE, line=LIGHT_BLUE)
    tb = slide.shapes.add_textbox(l + Inches(0.12), t + Inches(0.08),
                                  w - Inches(0.24), h - Inches(0.16))
    tf = tb.text_frame; tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line; p.font.size = Pt(sz)
        p.font.color.rgb = CODE_GREEN; p.font.name = "Consolas"
        p.space_after = Pt(1)

def _divider(slide, top):
    _shape(slide, MARGIN, top, LINE_W, Inches(0.04), fill=GOLD)

def _sub_heading(slide, l, t, w, text, sz=20):
    _tb(slide, l, t, w, Inches(0.45), text, sz=sz, color=ORANGE, bold=True)

def _add_swirl(slide):
    if os.path.exists(SWIRL_PATH):
        slide.shapes.add_picture(SWIRL_PATH, Inches(4.71), Inches(3.29),
                                 Inches(8.81), Inches(4.21))

# ── Delete all slides ───────────────────────────────────────────
def _delete_all_slides(prs):
    sldIdLst = prs.slides._sldIdLst
    for sldId in list(sldIdLst):
        rId = sldId.get(qn('r:id'))
        prs.part.drop_rel(rId)
        sldIdLst.remove(sldId)

prs = Presentation(TEMPLATE)
blank_layout = prs.slide_layouts[21]  # 2_Blank-White
_delete_all_slides(prs)

slide_idx = 0

def new_slide(bg=DARK_BG, add_swirl=False):
    global slide_idx
    slide_idx += 1
    s = prs.slides.add_slide(blank_layout)
    _set_bg(s, bg)
    if add_swirl:
        _add_swirl(s)
    return s

# ═══════════════════════════════════════════════════════════════
#  1  TITLE
# ═══════════════════════════════════════════════════════════════
s = new_slide(add_swirl=True)
_divider(s, Inches(3.0))
_tb(s, Inches(1), Inches(1.3), Inches(11), Inches(1.2),
    "Neo4j Fondamentali", sz=46, color=GOLD, bold=True, align=PP_ALIGN.CENTER)
_tb(s, Inches(1), Inches(3.3), Inches(11), Inches(0.8),
    "Sessione 1: Primi Passi con i Grafi", sz=26, color=WHITE, align=PP_ALIGN.CENTER)
_tb(s, Inches(1), Inches(4.8), Inches(11), Inches(0.5),
    "Database a grafo, Cypher e l\u2019ecosistema Neo4j", sz=17, color=LIGHT_GRAY,
    align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
#  2  AGENDA  (2 colonne)
# ═══════════════════════════════════════════════════════════════
s = new_slide()
_title_heading(s, "Agenda")
_tb(s, MARGIN, Inches(1.3), CONTENT_W, Inches(0.4),
    "Teoria", sz=18, color=ORANGE, bold=True)
_bullets(s, MARGIN, Inches(1.8), Inches(4.2), Inches(3.5), [
    "Cosa sono i database a grafo",
    "Property Graph Model",
    "Ecosistema Neo4j",
    "Connettersi con Python",
    "Cypher: pattern matching",
], sz=15, color=WHITE, spacing=1.3)
_tb(s, Inches(5.2), Inches(1.3), Inches(4.2), Inches(0.4),
    "Pratica", sz=18, color=ORANGE, bold=True)
_bullets(s, Inches(5.2), Inches(1.8), Inches(4.2), Inches(3.5), [
    "Cypher: leggere e scrivere dati",
    "Vincoli e integrit\u00e0",
    "Strategie di caricamento",
    "Indici e performance",
    "Visualizzare il grafo",
], sz=15, color=WHITE, spacing=1.3)

# ═══════════════════════════════════════════════════════════════
#  3  COSA SONO I DATABASE A GRAFO
# ═══════════════════════════════════════════════════════════════
s = new_slide()
_title_heading(s, "Cosa sono i database a grafo?")
_bullets(s, MARGIN, Inches(1.25), CONTENT_W, Inches(2.3), [
    "Rappresentano dati come nodi (entit\u00e0) e relazioni (connessioni)",
    "Le relazioni sono cittadini di prima classe \u2014 non calcolate via JOIN",
    "Ogni nodo/relazione pu\u00f2 avere propriet\u00e0 (coppie chiave-valore)",
    "Ideali per: gerarchie, reti sociali, raccomandazioni, frodi, grafi della conoscenza",
], sz=16, color=WHITE, spacing=1.3)
_divider(s, Inches(3.8))
_sub_heading(s, MARGIN, Inches(4.05), CONTENT_W, "Differenza chiave: DB relazionali vs grafo")
_bullets(s, MARGIN, Inches(4.6), CONTENT_W, Inches(2.5), [
    "SQL: JOIN su tabelle \u2192 costo cresce con profondit\u00e0 delle relazioni",
    "Grafo: le relazioni sono puntatori diretti \u2192 traversal in O(1)",
    "Esempio: \u201ctrova cosa hanno in comune 2 persone\u201d \u2192 pochi hop nel grafo",
], sz=14, color=LIGHT_GRAY, spacing=1.2)

# ═══════════════════════════════════════════════════════════════
#  4  PROPERTY GRAPH MODEL
# ═══════════════════════════════════════════════════════════════
s = new_slide()
_title_heading(s, "Il Property Graph Model")
_sub_heading(s, MARGIN, Inches(1.2), Inches(4.2), "Nodi")
_bullets(s, MARGIN, Inches(1.7), Inches(4.2), Inches(1.8), [
    "Entit\u00e0: persone, progetti, skills, luoghi",
    "Etichette (labels): :Person, :Project, :Skill",
    "Propriet\u00e0: coppie chiave-valore",
    "Un nodo pu\u00f2 avere pi\u00f9 labels",
], sz=14, color=WHITE, spacing=1.2)
_sub_heading(s, Inches(5.2), Inches(1.2), Inches(4.2), "Relazioni")
_bullets(s, Inches(5.2), Inches(1.7), Inches(4.2), Inches(1.8), [
    "Collegano due nodi con direzione",
    "Tipo: HAS_SKILL, WORKED_ON, REPORTS_TO",
    "Possono avere propriet\u00e0 (proficiency, anno)",
    "Sempre un inizio e una fine (grafo diretto)",
], sz=14, color=WHITE, spacing=1.2)
_divider(s, Inches(3.9))
_tb(s, MARGIN, Inches(4.1), CONTENT_W, Inches(0.4),
    "Pattern Cypher (notazione ASCII-art):", sz=16, color=GOLD, bold=True)
_code(s, MARGIN, Inches(4.65), Inches(9.0), Inches(2.2), [
    '(:Person {nome: "Alice"})            // nodo con label e propriet\u00e0',
    '    |',
    '  [:HAS_SKILL {livello: "Senior"}]   // relazione con propriet\u00e0',
    '    \u2514\u2500\u25b6 (:Skill {nome: "Python"})        // nodo destinazione',
], sz=12)

# ═══════════════════════════════════════════════════════════════
#  5  PERCHÉ NEO4J
# ═══════════════════════════════════════════════════════════════
s = new_slide()
_title_heading(s, "Perch\u00e9 Neo4j?")
_bullets(s, MARGIN, Inches(1.25), CONTENT_W, Inches(5.5), [
    "Native graph storage \u2014 architettura ottimizzata per grafi",
    "Cypher: linguaggio dichiarativo pattern-based (standard de-facto)",
    "ACID compliant \u2014 transazioni complete con rollback",
    "Index-free adjacency \u2014 traversal senza lookup di indice",
    "Ecosistema ricco: AuraDB, Browser, Bloom, Workspace, Drivers",
    "Open source (Community) + Enterprise (clustering, backup)",
    "OPENCypher \u2192 standard ISO (GQL in arrivo)",
], sz=15, color=WHITE, spacing=1.3)

# ═══════════════════════════════════════════════════════════════
#  6  ECOSISTEMA NEO4J
# ═══════════════════════════════════════════════════════════════
s = new_slide()
_title_heading(s, "L\u2019Ecosistema Neo4j")
_bullets(s, MARGIN, Inches(1.25), CONTENT_W, Inches(5), [
    "Neo4j Browser \u2014 UI web per Cypher (localhost:7474)",
    "Neo4j AuraDB \u2014 cloud managed (Free Tier: 200k nodi)",
    "Neo4j Bloom \u2014 visualizzazione senza scrivere Cypher",
    "Neo4j Workspace \u2014 modellazione collaborativa",
    "Neo4j Desktop \u2014 dev environment locale",
    "Drivers: Python, Java, JavaScript, Go, .NET, C",
], sz=15, color=WHITE, spacing=1.3)

# ═══════════════════════════════════════════════════════════════
#  7  CONNESSIONE PYTHON
# ═══════════════════════════════════════════════════════════════
s = new_slide()
_title_heading(s, "Connettersi a Neo4j con Python")
_bullets(s, MARGIN, Inches(1.15), CONTENT_W, Inches(1.5), [
    "Driver ufficiale: pip install neo4j",
    "Pattern: GraphDatabase.driver(URI, auth=(user, password))",
    "Driver mantiene un pool di connessioni (thread-safe)",
    "Usare context manager o driver.close()",
], sz=15, color=WHITE, spacing=1.2)
_code(s, MARGIN, Inches(3.1), CODE_W, Inches(3.5), [
    'from neo4j import GraphDatabase',
    'import os',
    '',
    'URI  = os.getenv("NEO4J_URI", "bolt://localhost:7687")',
    'USER = os.getenv("NEO4J_USERNAME", "neo4j")',
    'PASS = os.getenv("NEO4J_PASSWORD", "password")',
    '',
    'driver = GraphDatabase.driver(URI, auth=(USER, PASS))',
    'driver.verify_connectivity()',
    '# "Connected to neo4j!"',
    '',
    'records, _, _ = driver.execute_query(',
    '    "MATCH (n) RETURN count(n) AS cnt")',
], sz=11)

# ═══════════════════════════════════════════════════════════════
#  8  CYPHER: PATTERN MATCHING
# ═══════════════════════════════════════════════════════════════
s = new_slide()
_title_heading(s, "Cypher: Pattern Matching")
_tb(s, MARGIN, Inches(1.1), CONTENT_W, Inches(0.4),
    "Cypher usa una sintassi ASCII-art per descrivere pattern nel grafo:", sz=14, color=LIGHT_GRAY)
_sub_heading(s, MARGIN, Inches(1.6), Inches(4.2), "Pattern per nodi", sz=18)
_code(s, MARGIN, Inches(2.1), Inches(4.4), Inches(1.4), [
    '(n)',
    '(:Person)',
    '(p:Person)',
    '(p:Person {nome:"Alice"})',
], sz=11)
_sub_heading(s, Inches(5.2), Inches(1.6), Inches(4.2), "Pattern per relazioni", sz=18)
_code(s, Inches(5.2), Inches(2.1), Inches(4.4), Inches(1.4), [
    '(p)-[:HAS_SKILL]->(s)',
    '(p)<-[:REPORTS_TO]-(s)',
    '(p)-[r:HAS_SKILL]->(s)',
    '(p)-[r {livello:"Senior"}]->(s)',
], sz=11)
_divider(s, Inches(3.85))
_tb(s, MARGIN, Inches(4.1), CONTENT_W, Inches(0.4),
    "Struttura base di una query:", sz=16, color=GOLD, bold=True)
_code(s, MARGIN, Inches(4.6), CODE_W, Inches(2.2), [
    'MATCH pattern   -- cerca nel grafo',
    'WHERE condizione -- filtra (opzionale)',
    'RETURN expr     -- proietta i risultati',
    '',
    'MATCH (p:Person)',
    'WHERE p.location = "Roma"',
    'RETURN p.nome, p.eta',
    'ORDER BY p.eta DESC LIMIT 10;',
], sz=11)

# ═══════════════════════════════════════════════════════════════
#  9  CYPHER: LEGGERE DATI
# ═══════════════════════════════════════════════════════════════
s = new_slide()
_title_heading(s, "Cypher: Leggere Dati")
_bullets(s, MARGIN, Inches(1.15), CONTENT_W, Inches(1.5), [
    "MATCH + RETURN: struttura portante di Cypher",
    "WHERE con AND, OR, IN, CONTAINS, =, >, <, STARTS WITH",
    "Aggregazioni: count(), collect(), avg(), sum(), min(), max()",
    "Ordinamento: ORDER BY ... ASC/DESC, LIMIT, SKIP",
], sz=14, color=WHITE, spacing=1.2)
_code(s, MARGIN, Inches(3.0), CODE_W, Inches(3.8), [
    'MATCH (p:Person)',
    'WHERE p.location = "Roma"',
    'RETURN p.nome, p.eta',
    'ORDER BY p.eta DESC LIMIT 10;',
    '',
    'MATCH (p:Person)-[:HAS_SKILL]->(s:Skill)',
    'RETURN s.nome, count(*) AS persone',
    'ORDER BY persone DESC;',
    '',
    'MATCH (p:Person)-[h:HAS_SKILL]->(s:Skill)',
    'RETURN s.nome, avg(h.livello) AS media',
    'ORDER BY media DESC;',
], sz=11)

# ═══════════════════════════════════════════════════════════════
#  10  CYPHER: SCRIVERE DATI
# ═══════════════════════════════════════════════════════════════
s = new_slide()
_title_heading(s, "Cypher: Scrivere Dati")
_bullets(s, MARGIN, Inches(1.15), CONTENT_W, Inches(1.5), [
    "CREATE \u2014 crea nodi/relazioni (fallisce se gi\u00e0 esiste)",
    "MERGE \u2014 cerca o crea (pattern match + insert if not exists)",
    "SET \u2014 aggiorna propriet\u00e0 o ne aggiunge di nuove",
    "DELETE + DETACH DELETE \u2014 rimuove nodi",
], sz=14, color=WHITE, spacing=1.2)
_code(s, MARGIN, Inches(3.0), Inches(9.0), Inches(3.5), [
    'CREATE (p:Person {id: 1, nome: "Alice", location: "Roma"});',
    '',
    'MERGE (p:Person {id: 1})',
    '  SET p.nome = "Alice", p.location = "Roma";',
    '',
    'MATCH (p:Person {id: 1}), (s:Skill {nome: "Python"})',
    'MERGE (p)-[h:HAS_SKILL]->(s)',
    '  SET h.livello = "Senior";',
], sz=11)

# ═══════════════════════════════════════════════════════════════
#  11  CONSTRAINTS
# ═══════════════════════════════════════════════════════════════
s = new_slide()
_title_heading(s, "Vincoli (Constraints)")
_bullets(s, MARGIN, Inches(1.15), CONTENT_W, Inches(2.2), [
    "UNIQUE \u2014 unicit\u00e0 su una propriet\u00e0 (es: Person.id)",
    "EXISTS \u2014 propriet\u00e0 obbligatoria sempre presente",
    "Node Key \u2014 unicit\u00e0 + exists combinati",
    "CREATE CONSTRAINT ... IF NOT EXISTS",
    "Creare PRIMA del caricamento",
], sz=14, color=WHITE, spacing=1.2)
_code(s, MARGIN, Inches(3.7), Inches(9.0), Inches(2.8), [
    'CREATE CONSTRAINT person_id IF NOT EXISTS',
    '  FOR (p:Person) REQUIRE p.id IS UNIQUE;',
    '',
    '-- Verificare i constraint:',
    'SHOW CONSTRAINTS;',
    '',
    'SHOW INDEXES;  -- anche gli indici creati',
], sz=11)

# ═══════════════════════════════════════════════════════════════
#  12  LOADING STRATEGIES
# ═══════════════════════════════════════════════════════════════
s = new_slide()
_title_heading(s, "Strategie di Caricamento")
_bullets(s, MARGIN, Inches(1.15), CONTENT_W, Inches(1.5), [
    "Row-by-row: una query per ogni riga \u2192 N chiamate",
    "UNWIND batch: colleziona dati, invia in un\u2019unica transazione",
    "MERGE evita duplicati, CREATE \u00e8 pi\u00f9 veloce ma rischioso",
    "Sempre batch con UNWIND + MERGE per dati nuovi",
], sz=14, color=WHITE, spacing=1.2)
_code(s, MARGIN, Inches(3.0), Inches(9.0), Inches(3.8), [
    '# Collezione dati (Python):',
    'batch = df[[\'id\',\'nome\',\'location\']].to_dict("records")',
    '',
    '# UNWIND bulk (Cypher):',
    'UNWIND $batch AS row',
    'MERGE (p:Person {id: row.id})',
    '  SET p.nome = row.nome, p.location = row.location',
    '',
    '# Performance:',
    '# 31.500 query singole  440 secondi',
    '# 7 query UNWIND         1.4 secondi  (300x)',
], sz=11)

# ═══════════════════════════════════════════════════════════════
#  13  INDICI
# ═══════════════════════════════════════════════════════════════
s = new_slide()
_title_heading(s, "Indici e Performance")
_bullets(s, MARGIN, Inches(1.15), CONTENT_W, Inches(2.2), [
    "UNIQUE crea automaticamente un indice",
    "Tipi: RANGE (default), TEXT, POINT (geospaziale)",
    "CREATE INDEX ... FOR (p:Person) ON (p.location)",
    "PROFILE / EXPLAIN per vedere se usa indici",
    "Indici migliorano LETTURA ma rallentano SCRITTURA",
], sz=14, color=WHITE, spacing=1.2)
_code(s, MARGIN, Inches(3.7), CODE_W, Inches(2.8), [
    'CREATE INDEX person_location IF NOT EXISTS',
    '  FOR (p:Person) ON (p.location);',
    '',
    'PROFILE',
    'MATCH (p:Person) WHERE p.location = "Roma" RETURN p;',
], sz=11)

# ═══════════════════════════════════════════════════════════════
#  14  VISUALIZZAZIONE
# ═══════════════════════════════════════════════════════════════
s = new_slide()
_title_heading(s, "Visualizzare il Grafo")
_bullets(s, MARGIN, Inches(1.15), CONTENT_W, Inches(2.5), [
    "Neo4j Browser \u2014 :play movie-graph per esempio interattivo",
    "Neo4j Bloom \u2014 drag & drop, storytelling visuale",
    "Python networkx + matplotlib \u2014 grafi statici",
    "Python pyvis \u2014 grafi interattivi HTML",
    "yFiles, Linkurious, Keylines \u2014 enterprise",
], sz=14, color=WHITE, spacing=1.2)
_code(s, MARGIN, Inches(4.0), CODE_W, Inches(2.5), [
    '# Grafo statico con networkx + matplotlib:',
    'import networkx as nx, matplotlib.pyplot as plt',
    'G = nx.Graph()',
    'G.add_edge("Alice", "Python", relation="HAS_SKILL")',
    'nx.draw(G, with_labels=True)',
    'plt.show()',
], sz=11)

# ═══════════════════════════════════════════════════════════════
#  15  MOVIE GRAPH
# ═══════════════════════════════════════════════════════════════
s = new_slide()
_title_heading(s, "Esempio: Movie Graph")
_bullets(s, MARGIN, Inches(1.15), CONTENT_W, Inches(1.3), [
    "Dataset incluso in Neo4j Browser: :play movie-graph",
    "Nodi: Movie, Person (attori/registi)",
    "Relazioni: ACTED_IN, DIRECTED, REVIEWED",
], sz=14, color=WHITE, spacing=1.2)
_code(s, MARGIN, Inches(2.8), Inches(9.0), Inches(4.0), [
    '// Creare il dataset: esegui dal Browser',
    ':play movie-graph',
    '',
    '// Attori di film usciti dopo il 2000:',
    'MATCH (p:Person)-[:ACTED_IN]->(m:Movie)',
    'WHERE m.released > 2000',
    'RETURN p.name, m.title, m.released',
    'ORDER BY m.released DESC;',
    '',
    '// Co-protagonisti di Tom Hanks:',
    'MATCH (tom:Person {name: "Tom Hanks"})',
    '  -[:ACTED_IN]->(m:Movie)<-[:ACTED_IN]-(co)',
    'RETURN co.name, m.title;',
], sz=11)

# ═══════════════════════════════════════════════════════════════
#  16  ESERCIZI
# ═══════════════════════════════════════════════════════════════
s = new_slide()
_title_heading(s, "Esercizi")
_bullets(s, MARGIN, Inches(1.25), CONTENT_W, Inches(4.5), [
    "Carica Movie Graph, trova attori di un film specifico",
    "Query: film diretti da un regista con anno e rating",
    "Trova attori che NON hanno recitato con Tom Hanks",
    "Connettiti via Python driver, esegui MATCH (n) RETURN n LIMIT 25",
    "Crea un nuovo nodo Person e una relazione ACTED_IN",
    "Esplora il dataset del progetto nel notebook",
], sz=15, color=WHITE, spacing=1.3)
_tb(s, MARGIN, Inches(5.8), CONTENT_W, Inches(0.4),
    "Le soluzioni sono nel notebook della Sessione 1", sz=13, color=GOLD)

# ═══════════════════════════════════════════════════════════════
#  17  TAKEAWAYS
# ═══════════════════════════════════════════════════════════════
s = new_slide()
_title_heading(s, "Punti Chiave")
_bullets(s, MARGIN, Inches(1.25), CONTENT_W, Inches(5.5), [
    "DB a grafo eccellono con dati connessi (no JOIN costosi)",
    "Neo4j: Property Graph Model \u2014 nodi + relazioni + labels",
    "Cypher \u00e8 dichiarativo: descrivi il pattern, non il percorso",
    "MERGE cerca o crea \u2014 caricamenti idempotenti",
    "UNWIND batching: 300\u00d7 pi\u00f9 veloce del row-by-row",
    "Constraints e indici vanno creati PRIMA del caricamento",
    "Neo4j Browser: esplorazione rapida e test query",
], sz=15, color=WHITE, spacing=1.3)

# ═══════════════════════════════════════════════════════════════
#  18  PROSSIMA SESSIONE
# ═══════════════════════════════════════════════════════════════
s = new_slide()
_divider(s, Inches(2.6))
_tb(s, Inches(1), Inches(0.8), Inches(11), Inches(1.2),
    "Alla prossima!", sz=42, color=GOLD, bold=True, align=PP_ALIGN.CENTER)
_tb(s, Inches(1), Inches(2.9), Inches(11), Inches(0.6),
    "Sessione 2: Cypher Avanzato", sz=24, color=WHITE, align=PP_ALIGN.CENTER)
_bullets(s, Inches(2.5), Inches(3.9), Inches(9), Inches(2.5), [
    "Pattern matching avanzato (path variables, shortestPath)",
    "Aggregazioni complesse e subquery",
    "Date, liste, map in Cypher",
    "Query tuning con PROFILE ed EXPLAIN",
    "APOC: procedure standard di Neo4j",
], sz=14, color=LIGHT_GRAY, spacing=1.3)

# ── Save ───────────────────────────────────────────────────────
OUT_PATH = os.path.join(ASSETS, 'session-01.pptx')
prs.save(OUT_PATH)
print(f"Salvato {OUT_PATH} ({len(prs.slides)} slide)")
